import streamlit as st
import os
import io
from pypdf import PdfReader
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import pandas as pd
import json
from queue import Queue
import threading
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse
import docx
import pptx
import time
import pickle
from typing import List, Tuple, Union
from docx import Document
from dotenv import load_dotenv
from groq import Groq
from pymongo import MongoClient
import sqlite3
from concurrent.futures import ThreadPoolExecutor

# Load environment variables
load_dotenv()


def connect_mongodb(uri: str):
    client = MongoClient(uri)
    return client

def connect_sqlite(db_path: str):
    conn = sqlite3.connect(db_path)
    return conn

@st.cache_resource
def get_embedding_model():
    return SentenceTransformer('all-MiniLM-L6-v2')

class DocumentProcessor:
    def __init__(self, mongo_uri=None, sqlite_path=None):
        self.vector_dimension = 384
        self.embedding_model = get_embedding_model()
        self.index = faiss.IndexFlatL2(self.vector_dimension)
        self.text_chunks = []
        self.chunk_sources = []
        self.processing_queue = Queue()
        self.processed_urls = set()
        self.mongo_client = connect_mongodb(mongo_uri) if mongo_uri else None
        self.sqlite_conn = connect_sqlite(sqlite_path) if sqlite_path else None

        # Initialize Groq client
        try:
            groq_api_key = os.getenv('GROQ_API_KEY')
            if not groq_api_key:
                raise ValueError("GROQ_API_KEY environment variable is not set")
            self.client = Groq(api_key=groq_api_key)
        except Exception as e:
            st.error(f"Error initializing Groq client: {str(e)}")
            self.client = None
    def query_mongodb(self, query: str):
        if not self.mongo_client:
            return []
        db = self.mongo_client['your_database']
        collection = db['your_collection']
        results = collection.find({"$text": {"$search": query}})
        return [(doc['text'], 'MongoDB') for doc in results]

    def query_sqlite(self, query: str):
        if not self.sqlite_conn:
            return []
        cursor = self.sqlite_conn.cursor()
        cursor.execute("SELECT text FROM your_table WHERE text LIKE ?", ('%' + query + '%',))
        results = cursor.fetchall()
        return [(row[0], 'SQLite') for row in results]

    def extract_text_from_pdf(self, file_content: bytes) -> str:
        try:
            with io.BytesIO(file_content) as pdf_file:
                pdf_reader = PdfReader(pdf_file)
                text = "".join(page.extract_text() for page in pdf_reader.pages)
            return text
        except Exception as e:
            st.error(f"Error processing PDF: {str(e)}")
            return ""

    def extract_text_from_docx(self, file_content: bytes) -> str:
        try:
            with io.BytesIO(file_content) as docx_file:
                doc = docx.Document(docx_file)
                text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
            return text
        except Exception as e:
            st.error(f"Error processing DOCX: {str(e)}")
            return ""

    def extract_text_from_pptx(self, file_content: bytes) -> str:
        try:
            with io.BytesIO(file_content) as pptx_file:
                prs = pptx.Presentation(pptx_file)
                text = "\n".join(shape.text for slide in prs.slides 
                               for shape in slide.shapes if hasattr(shape, "text"))
            return text
        except Exception as e:
            st.error(f"Error processing PPTX: {str(e)}")
            return ""

    def extract_text_from_txt(self, file_content: bytes) -> str:
        try:
            return file_content.decode('utf-8')
        except Exception as e:
            st.error(f"Error processing TXT: {str(e)}")
            return ""

    def extract_data_from_excel(self, file_content: bytes) -> pd.DataFrame:
        try:
            with io.BytesIO(file_content) as excel_file:
                df = pd.read_excel(excel_file)
            return df
        except Exception as e:
            st.error(f"Error processing Excel: {str(e)}")
            return pd.DataFrame()

    def extract_data_from_csv(self, file_content: bytes) -> pd.DataFrame:
        try:
            with io.BytesIO(file_content) as csv_file:
                df = pd.read_csv(csv_file)
            return df
        except Exception as e:
            st.error(f"Error processing CSV: {str(e)}")
            return pd.DataFrame()

    def extract_data_from_json(self, file_content: bytes) -> pd.DataFrame:
        try:
            data = json.loads(file_content.decode('utf-8'))
            df = pd.json_normalize(data)
            return df
        except Exception as e:
            st.error(f"Error processing JSON: {str(e)}")
            return pd.DataFrame()

    def process_file(self, file_content: bytes, file_name: str) -> Union[str, pd.DataFrame]:
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        try:
            extension = file_name.lower().split('.')[-1]
            status_text.text(f"Processing {file_name}...")
            progress_bar.progress(25)
            
            # Use the existing extractors dictionary
            extractors = {
                'pdf': self.extract_text_from_pdf,
                'docx': self.extract_text_from_docx,
                'pptx': self.extract_text_from_pptx,
                'txt': self.extract_text_from_txt,
                'xlsx': self.extract_data_from_excel,
                'csv': self.extract_data_from_csv,
                'json': self.extract_data_from_json
            }
            
            extractor = extractors.get(extension)
            if extractor:
                result = extractor(file_content)
                progress_bar.progress(100)
                status_text.text(f"Completed processing {file_name}")
                time.sleep(0.5)
                status_text.empty()
                progress_bar.empty()
                return result
            else:
                progress_bar.empty()
                status_text.empty()
                st.error(f"Unsupported file type: {extension}")
                return None
                
        except Exception as e:
            progress_bar.empty()
            status_text.empty()
            st.error(f"Error processing {file_name}: {str(e)}")
            return None


    def split_text(self, text: str, chunk_size: int = 1000) -> List[str]:
        if not text.strip():
            return []
        words = text.split()
        chunks = []
        current_chunk = []
        current_size = 0
        
        for word in words:
            current_chunk.append(word)
            current_size += len(word) + 1
            
            if current_size >= chunk_size:
                chunks.append(" ".join(current_chunk))
                current_chunk = []
                current_size = 0
        
        if current_chunk:
            chunks.append(" ".join(current_chunk))
        
        return chunks

    def delayed_process_documents(self):
        while True:
            if not self.processing_queue.empty():
                content, source_name = self.processing_queue.get()
                time.sleep(30)  # 30-second delay
                self.process_document(content, source_name)
                self.save_index()
            time.sleep(1)

    def start_processing_thread(self):
        thread = threading.Thread(target=self.delayed_process_documents, daemon=True)
        thread.start()

    def process_document(self, content: Union[str, pd.DataFrame], source_name: str) -> int:
        try:
            if isinstance(content, pd.DataFrame):
                # Optimize DataFrame processing
                content = content.head(1000).to_string()  # Limit initial processing
            
            # Use smaller chunk size for faster processing
            chunks = self.split_text(content, chunk_size=500)
            
            if not chunks:
                return 0
                
            # Batch process embeddings for speed
            batch_size = 32
            all_embeddings = []
            for i in range(0, len(chunks), batch_size):
                batch = chunks[i:i+batch_size]
                embeddings = self.embedding_model.encode(batch, batch_size=batch_size)
                all_embeddings.append(embeddings)
            
            embeddings = np.vstack(all_embeddings)
            faiss.normalize_L2(embeddings)
            
            self.index.add(embeddings.astype('float32'))
            self.text_chunks.extend(chunks)
            self.chunk_sources.extend([source_name] * len(chunks))
            
            return len(chunks)
        except Exception as e:
            st.error(f"Error processing document {source_name}: {str(e)}")
            return 0

    def save_index(self, directory: str = "./saved_index") -> None:
        try:
            os.makedirs(directory, exist_ok=True)
            faiss.write_index(self.index, os.path.join(directory, "docs.index"))
            with open(os.path.join(directory, "chunks.pkl"), "wb") as f:
                pickle.dump((self.text_chunks, self.chunk_sources), f)
        except Exception as e:
            st.error(f"Error saving index: {str(e)}")

    def load_index(self, directory: str = "./saved_index") -> bool:
        try:
            index_path = os.path.join(directory, "docs.index")
            chunks_path = os.path.join(directory, "chunks.pkl")
            
            if os.path.exists(index_path) and os.path.exists(chunks_path):
                self.index = faiss.read_index(index_path)
                with open(chunks_path, "rb") as f:
                    self.text_chunks, self.chunk_sources = pickle.load(f)
                return True
            return False
        except Exception as e:
            st.error(f"Error loading index: {str(e)}")
            return False

    def query_documents(self, query: str, n_results: int = 3) -> Tuple[str, List[Tuple[str, str]], bool]:
        try:
            # Cache query embeddings
            @st.cache_data(ttl=3600)
            def get_query_embedding(q):
                return self.embedding_model.encode([q])

            query_embedding = get_query_embedding(query)
            faiss.normalize_L2(query_embedding)
            
            # Get FAISS results
            distances, indices = self.index.search(query_embedding.astype('float32'), n_results)
            relevant_chunks = [(self.text_chunks[i], self.chunk_sources[i]) for i in indices[0]]
            
            # Get database results sequentially
            mongo_results = self.query_mongodb(query) if self.mongo_client else []
            sqlite_results = self.query_sqlite(query) if self.sqlite_conn else []

            # Combine results
            all_results = relevant_chunks + mongo_results + sqlite_results

            # Optimize context building
            context = "\n".join(chunk for chunk, _ in all_results[:5])  # Limit context size
            
            if not self.client:
                return "Error: Groq client not initialized properly.", all_results, False

            # Optimize prompt for faster response
            prompt = f'''Given this context: {context}
    Query: {query}
    Provide a concise, factual answer based on the context.'''
            
            completion = self.client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama-3.1-8b-instant",
                temperature=0.5,  # Lower temperature for more focused responses
            )
            return completion.choices[0].message.content, all_results, False

        except Exception as e:
            st.error(f"Error: {str(e)}")
            return "Error processing query. Please try again.", [], False


    def crawl_website(self, base_url: str, max_pages: int = 50) -> List[Tuple[str, str]]:
        pages_content = []
        to_visit = {base_url}
        visited = set()

        while to_visit and len(visited) < max_pages:
            url = to_visit.pop()
            if url in visited:
                continue

            try:
                headers = {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                }
                response = requests.get(url, headers=headers, timeout=10)
                response.raise_for_status()
                
                soup = BeautifulSoup(response.text, 'html.parser')
                for script in soup(["script", "style", "meta", "noscript"]):
                    script.decompose()
                text = soup.get_text()
                pages_content.append((url, text))

                links = soup.find_all('a', href=True)
                for link in links:
                    href = urljoin(url, link['href'])
                    if self.is_same_domain(url, href) and href not in visited:
                        to_visit.add(href)
                
                visited.add(url)
                time.sleep(1)
                
            except Exception as e:
                st.error(f"Error crawling {url}: {str(e)}")
                continue

        return pages_content

    def is_same_domain(self, url1: str, url2: str) -> bool:
        domain1 = urlparse(url1).netloc
        domain2 = urlparse(url2).netloc
        return domain1 == domain2

def main():
    
    st.set_page_config(
        page_title="Quick RAG Search",
        page_icon="🔍",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    # Add custom CSS for better styling
    st.markdown("""
        <style>
        .stApp {
            max-width: 1200px;
            margin: 0 auto;
        }
        .stButton>button {
            width: 100%;
        }
        .stProgress .st-bo {
            background-color: #00ff00;
        }
        </style>
    """, unsafe_allow_html=True)
    
    # Create two columns for better layout
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.title("🔍 Quick RAG Search")
        st.markdown("Upload documents, ask questions, get instant answers!")
    
    # Check for GROQ_API_KEY
    if not os.getenv('GROQ_API_KEY'):
        st.error("Please set the GROQ_API_KEY environment variable")
        return

    if 'processor' not in st.session_state:
        st.session_state.processor = DocumentProcessor()
        st.session_state.processor.start_processing_thread()
    
    if not st.session_state.get("index_loaded"):
        if st.session_state.processor.load_index():
            st.session_state.index_loaded = True
            st.success("Loaded existing knowledge base!")
        else:
            st.info("No existing knowledge base found. Please upload documents or add URLs.")
    with st.sidebar:
        st.header("Database Connections")
        
        # MongoDB Connection
        mongo_uri = st.text_input("MongoDB URI", key="mongo_uri")
        if st.button("Connect to MongoDB"):
            if mongo_uri:
                st.session_state.mongo_client = connect_mongodb(mongo_uri)
                st.success("Connected to MongoDB!")
                st.session_state.processor = DocumentProcessor(mongo_uri=mongo_uri, sqlite_path=sqlite_path)
                st.session_state.processor.start_processing_thread()
            else:
                st.error("Please provide a valid MongoDB URI.")
        
        # SQLite Connection
        sqlite_path = st.text_input("SQLite DB Path", key="sqlite_path")
        if st.button("Connect to SQLite"):
            if sqlite_path:
                st.session_state.sqlite_conn = connect_sqlite(sqlite_path)
                st.success("Connected to SQLite!")
                st.session_state.processor = DocumentProcessor(mongo_uri=mongo_uri, sqlite_path=sqlite_path)
                st.session_state.processor.start_processing_thread()
            else:
                st.error("Please provide a valid SQLite database path.")
        
        # Existing document and URL upload logic
        st.header("Add Documents")
        
        supported_types = ["pdf", "docx", "txt", "pptx", "xlsx", "csv", "json"]
        uploaded_files = st.file_uploader(
            f"Upload documents ({', '.join(supported_types)})", 
            type=supported_types, 
            accept_multiple_files=True
        )
        
        st.header("Add URLs")
        url_input = st.text_area("Enter URLs (one per line)")
        
        crawl_depth = st.slider("Maximum pages to crawl per website", 1, 100, 50)
        
        process_button = st.button("Process Documents")
        
        if process_button:
            total_queued = 0
            
            if uploaded_files:
                for uploaded_file in uploaded_files:
                    with st.spinner(f"Queueing {uploaded_file.name}..."):
                        file_content = uploaded_file.read()
                        content = st.session_state.processor.process_file(file_content, uploaded_file.name)
                        if isinstance(content, pd.DataFrame):
                            st.write(content)
                        elif content:
                            st.session_state.processor.processing_queue.put((content, uploaded_file.name))
                            total_queued += 1
            
            if url_input:
                urls = [url.strip() for url in url_input.split('\n') if url.strip()]
                for url in urls:
                    if url:
                        with st.spinner(f"Crawling {url}..."):
                            pages_content = st.session_state.processor.crawl_website(url, max_pages=crawl_depth)
                            for page_url, content in pages_content:
                                if content:
                                    st.session_state.processor.processing_queue.put((content, page_url))
                                    total_queued += 1
            
            if total_queued > 0:
                st.success(f"Queued {total_queued} documents for processing. They will be processed with a 30-second delay.")
            else:
                st.warning("No new content was queued. Please check your inputs.")

        if st.button("Clear Knowledge Base"):
            if os.path.exists("./saved_index"):
                try:
                    os.remove("./saved_index/docs.index")
                    os.remove("./saved_index/chunks.pkl")
                    st.session_state.processor = DocumentProcessor()
                    st.session_state.processor.start_processing_thread()
                    st.session_state.index_loaded = False
                    st.success("Knowledge base cleared successfully!")
                except Exception as e:
                    st.error(f"Error clearing knowledge base: {str(e)}")

    st.header("Ask Questions")
    query = st.text_input("Enter your question:")

    if query:
        if len(st.session_state.processor.text_chunks) == 0:
            st.warning("Please add some documents to the knowledge base first!")
        else:
            with st.spinner("Searching and generating response..."):
                answer, sources, _ = st.session_state.processor.query_documents(query)
                st.markdown("### Answer")
                st.write(answer)

    st.header("Generate Report")
    report_button = st.button("Generate Report")

    if report_button:
        if len(st.session_state.processor.text_chunks) == 0:
            st.warning("Please add some documents to the knowledge base first!")
        else:
            with st.spinner("Generating report..."):
                report_content = f"Answer:\n{answer}\n\nSources:\n" + "\n".join([f"{i+1}. {source}" for i, (_, source) in enumerate(sources)])
                
                doc = Document()
                doc.add_heading('Report', level=1)
                doc.add_paragraph(report_content)
                
                doc_path = "report.docx"
                doc.save(doc_path)
                
                st.success("Report generated successfully!")
                st.download_button(
                    label="Download Report",
                    data=open(doc_path, "rb").read(),
                    file_name="report.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )

if __name__ == "__main__":
    main()